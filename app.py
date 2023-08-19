from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
from flask import Flask, request, render_template
import os
import docx2txt
from PyPDF2 import PdfReader
from pptx import Presentation

app = Flask(__name__)

@app.route('/')
def index():
    return render_template('comparar.html')

@app.route('/procesar', methods=['POST'])
def procesar():
    resultado = None  # Variable para almacenar el resultado

    if 'archivo1' not in request.files or 'archivo2' not in request.files:
        return render_template('comparar.html', resultado=resultado)

    archivo1 = request.files['archivo1']
    archivo2 = request.files['archivo2']

    if archivo1.filename == '' or archivo2.filename == '':
        return render_template('comparar.html', resultado=resultado)

    if archivo1 and archivo2:
        # Guarda los archivos cargados temporalmente
        archivo1_path = os.path.join("uploads", archivo1.filename)
        archivo2_path = os.path.join("uploads", archivo2.filename)
        archivo1.save(archivo1_path)
        archivo2.save(archivo2_path)

        # Extraer texto de archivos PDF y PPT
        texto1 = extraer_texto_archivo(archivo1_path)
        texto2 = extraer_texto_archivo(archivo2_path)

        # Tokeniza el texto
        tokens1 = texto1.split()
        tokens2 = texto2.split()

        # Calcula la frecuencia inversa (TF-IDF) de los tokens
        vectorizer = TfidfVectorizer()
        tfidf_matrix = vectorizer.fit_transform([texto1, texto2])
        similarity_score = cosine_similarity(tfidf_matrix[0], tfidf_matrix[1])[0][0]

        # Mandar los archivos alv
        os.remove(archivo1_path)
        os.remove(archivo2_path)

        # Formatea el resultado como porcentaje
        similarity_score_percent = round(similarity_score * 100, 2)
        resultado = f"Similitud de los Archivos: {similarity_score_percent}%"

        return render_template('comparar.html', resultado=resultado)

    return render_template('comparar.html', resultado=resultado)

def extraer_texto_archivo(archivo_path):
    _, archivo_extension = os.path.splitext(archivo_path)
    archivo_extension = archivo_extension.lower()

    if archivo_extension == '.txt':
        with open(archivo_path, 'r', encoding='utf-8') as archivo:
            return archivo.read()
    elif archivo_extension == '.docx':
        return docx2txt.process(archivo_path)
    elif archivo_extension == '.pdf':
        return extraer_texto_pdf(archivo_path)
    elif archivo_extension == '.pptx':
        return extraer_texto_ppt(archivo_path)
    else:
        return ''

def extraer_texto_pdf(pdf_path):
    texto = ''
    pdf_reader = PdfReader(pdf_path)
    for page in pdf_reader.pages:
        texto += page.extract_text()
    return texto

def extraer_texto_ppt(ppt_path):
    texto = ''
    presentation = Presentation(ppt_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texto += shape.text
    return texto

if __name__ == '__main__':
    app.run(debug=True)
